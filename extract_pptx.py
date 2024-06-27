from pptx import Presentation
from deep_translator import GoogleTranslator

def translate_pptx(pptx_path, output_path, source_lang='auto', target_lang='vi'):
    prs = Presentation(pptx_path)
    translator = GoogleTranslator(source=source_lang, target=target_lang)

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    translated_text = translator.translate(run.text)
                    run.text = run.text + "\n" + translated_text

    prs.save(os.path.join(output_path, 'translate.pptx'))
 